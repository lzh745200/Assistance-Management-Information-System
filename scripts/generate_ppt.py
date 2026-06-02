"""Generate beautiful PPT presentation for the military rural revitalization system."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# Color scheme - military green theme
DARK_GREEN = RGBColor(0x1B, 0x3A, 0x2D)
MID_GREEN = RGBColor(0x2D, 0x5A, 0x3F)
LIGHT_GREEN = RGBColor(0x4A, 0x8C, 0x5C)
ACCENT_GOLD = RGBColor(0xC8, 0x96, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF0, 0xF5, 0xF1)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x1A)
GRAY_TEXT = RGBColor(0x66, 0x66, 0x66)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_multi_text(slide, left, top, width, height, lines, font_size=16,
                   color=DARK_TEXT, line_spacing=1.5):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Microsoft YaHei'
        p.space_after = Pt(font_size * (line_spacing - 1) * 2)
    return tf

def add_accent_bar(slide, left, top, width, height):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_GOLD
    shape.line.fill.background()
    return shape

# ============================================================
# SLIDE 1: TITLE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK_GREEN)

# Decorative line at top
add_accent_bar(slide, 0, 0, 13.333, 0.08)

# Main title
add_text_box(slide, 1.5, 1.8, 10.3, 1.5,
             "军队乡村振兴管理系统",
             font_size=52, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Subtitle
add_text_box(slide, 1.5, 3.3, 10.3, 0.8,
             "Military Rural Revitalization Management System",
             font_size=24, color=ACCENT_GOLD, alignment=PP_ALIGN.CENTER)

# Separator line
add_accent_bar(slide, 4.5, 4.3, 4.3, 0.04)

# Version and date
add_text_box(slide, 1.5, 4.7, 10.3, 0.6,
             "版本 v1.1.0  |  2026年4月  |  生产就绪",
             font_size=20, color=RGBColor(0xAA, 0xBB, 0xAA), alignment=PP_ALIGN.CENTER)

# Bottom description
add_text_box(slide, 1.5, 5.8, 10.3, 0.6,
             "完全离线的单机版桌面应用 · 数据安全 · 高效便捷",
             font_size=16, color=RGBColor(0x88, 0x99, 0x88), alignment=PP_ALIGN.CENTER)

# Bottom decorative bar
add_accent_bar(slide, 0, 7.42, 13.333, 0.08)

# ============================================================
# SLIDE 2: PROJECT OVERVIEW
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_accent_bar(slide, 0, 0, 13.333, 0.06)

add_text_box(slide, 0.8, 0.4, 5, 0.7, "项目概览", font_size=36, color=DARK_GREEN, bold=True)
add_accent_bar(slide, 0.8, 1.05, 2.5, 0.04)

# Left column - description
overview_lines = [
    "军队乡村振兴管理系统是一款专为军队参与乡村振兴工作设计的",
    "完全离线单机版桌面应用。",
    "",
    "系统覆盖帮扶村管理、项目管理、资金管理、政策法规、",
    "数据统计分析等全流程业务，支持多级组织架构和角色权限管理。",
    "",
    "采用 Electron 桌面壳 + FastAPI 后端 + Vue 3 前端架构，",
    "确保数据安全和操作便捷性，适用于军队内部网络环境。",
]
add_multi_text(slide, 0.8, 1.4, 7.5, 4.5, overview_lines, font_size=17, color=DARK_TEXT)

# Right column - key metrics (card style)
cards = [
    ("v1.1.0", "当前版本"),
    ("2812+", "后端测试用例"),
    ("1449+", "前端测试用例"),
    ("100/100", "系统健康评分"),
]
for i, (value, label) in enumerate(cards):
    x = 9.0
    y = 1.3 + i * 1.45
    # Card background
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(3.5), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BG
    shape.line.fill.background()
    shape.shadow.inherit = False

    add_text_box(slide, x + 0.2, y + 0.15, 3.1, 0.5,
                 value, font_size=28, color=DARK_GREEN, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.2, y + 0.65, 3.1, 0.4,
                 label, font_size=14, color=GRAY_TEXT, alignment=PP_ALIGN.CENTER)

add_accent_bar(slide, 0, 7.44, 13.333, 0.06)

# ============================================================
# SLIDE 3: TECHNICAL ARCHITECTURE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_accent_bar(slide, 0, 0, 13.333, 0.06)

add_text_box(slide, 0.8, 0.4, 5, 0.7, "技术架构", font_size=36, color=DARK_GREEN, bold=True)
add_accent_bar(slide, 0.8, 1.05, 2, 0.04)

# Architecture layers
layers = [
    ("桌面壳层", "Electron 28", "窗口管理 · 系统托盘 · 自动更新 · 进程守护", DARK_GREEN),
    ("前端层", "Vue 3 + TypeScript + Element Plus", "响应式UI · 组合式API · Pinia状态管理 · ECharts可视化", MID_GREEN),
    ("后端层", "Python 3.11 + FastAPI + SQLAlchemy 2.0", "RESTful API · DDD架构 · JWT认证 · 数据验证", LIGHT_GREEN),
    ("数据层", "SQLite + diskcache", "本地数据库 · 磁盘缓存 · Alembic迁移 · 自动备份", RGBColor(0x6B, 0xA3, 0x7A)),
]

for i, (layer_name, tech, desc, color) in enumerate(layers):
    y = 1.5 + i * 1.35
    # Layer block
    shape = slide.shapes.add_shape(1, Inches(0.8), Inches(y), Inches(11.7), Inches(1.15))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

    add_text_box(slide, 1.1, y + 0.1, 2.2, 0.4,
                 layer_name, font_size=18, color=WHITE, bold=True)
    add_text_box(slide, 1.1, y + 0.5, 2.2, 0.5,
                 tech, font_size=12, color=RGBColor(0xDD, 0xEE, 0xDD))
    add_text_box(slide, 3.6, y + 0.15, 8.5, 0.8,
                 desc, font_size=15, color=WHITE)

# Bottom features
features = [
    "离线优先", "DDD架构", "JWT认证", "RBAC权限",
    "自动迁移", "审计日志", "数据加密", "批量操作"
]
for i, feat in enumerate(features):
    x = 0.8 + i * 1.55
    shape = slide.shapes.add_shape(1, Inches(x), Inches(7.0), Inches(1.4), Inches(0.4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BG
    shape.line.fill.background()
    add_text_box(slide, x, 7.0, 1.4, 0.4,
                 feat, font_size=11, color=DARK_GREEN, bold=True, alignment=PP_ALIGN.CENTER)

add_accent_bar(slide, 0, 7.44, 13.333, 0.06)

# ============================================================
# SLIDE 4: CORE FEATURES
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_accent_bar(slide, 0, 0, 13.333, 0.06)

add_text_box(slide, 0.8, 0.4, 5, 0.7, "核心功能模块", font_size=36, color=DARK_GREEN, bold=True)
add_accent_bar(slide, 0.8, 1.05, 3, 0.04)

modules = [
    ("帮扶村管理", "村庄基础信息、年度数据\n跟踪、地理信息管理\n支持批量导入导出", "🏘️"),
    ("项目管理", "帮扶项目全生命周期管理\n进度跟踪、资金使用监控\n项目评估和审计", "📋"),
    ("资金管理", "帮扶资金台账\n多级审批工作流\n资金使用分析和报表", "💰"),
    ("组织管理", "军队单位组织架构\n多级层级管理\n数据权限隔离", "🏛️"),
    ("政策法规", "政策文件管理\n分类检索和全文搜索\n政策落实跟踪", "📜"),
    ("数据统计", "多维度数据分析\n可视化图表展示\n自定义报表生成", "📊"),
    ("系统管理", "用户和角色权限管理\n菜单权限配置\n审计日志和备份恢复", "⚙️"),
    ("数据同步", "增量数据包导入导出\n冲突检测和处理\n离线地图瓦片管理", "🔄"),
]

for i, (name, desc, icon) in enumerate(modules):
    col = i % 4
    row = i // 4
    x = 0.7 + col * 3.1
    y = 1.4 + row * 3.0

    # Card bg
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(2.9), Inches(2.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BG
    shape.line.fill.background()

    add_text_box(slide, x + 0.15, y + 0.15, 2.6, 0.5,
                 f"{icon}  {name}", font_size=18, color=DARK_GREEN, bold=True)
    add_text_box(slide, x + 0.2, y + 0.75, 2.5, 1.7,
                 desc, font_size=13, color=DARK_TEXT)

add_accent_bar(slide, 0, 7.44, 13.333, 0.06)

# ============================================================
# SLIDE 5: SYSTEM STATUS & TESTING
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_accent_bar(slide, 0, 0, 13.333, 0.06)

add_text_box(slide, 0.8, 0.4, 6, 0.7, "系统质量与健康状态", font_size=36, color=DARK_GREEN, bold=True)
add_accent_bar(slide, 0.8, 1.05, 3.5, 0.04)

# Quality metrics - large cards
metrics = [
    ("2812 / 2812", "后端测试通过", "pytest · 覆盖率 > 80%"),
    ("1449 / 1449", "前端测试通过", "Vitest · TypeScript 0 错误"),
    ("0", "代码质量问题", "Flake8 0 · ESLint 0"),
    ("100/100", "系统健康评分", "所有模块正常运行"),
]

for i, (value, label, detail) in enumerate(metrics):
    x = 0.6 + i * 3.15
    y = 1.5
    # Card
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(2.95), Inches(2.3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BG
    shape.line.fill.background()

    add_text_box(slide, x + 0.15, y + 0.3, 2.65, 0.7,
                 value, font_size=38, color=DARK_GREEN, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.15, y + 1.1, 2.65, 0.45,
                 label, font_size=18, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.15, y + 1.6, 2.65, 0.4,
                 detail, font_size=12, color=GRAY_TEXT, alignment=PP_ALIGN.CENTER)

# Database tables
add_text_box(slide, 0.8, 4.3, 6, 0.5, "核心数据表", font_size=24, color=DARK_GREEN, bold=True)

tables = [
    "users (用户)", "villages (帮扶村)", "projects (项目)", "organizations (组织)",
    "funds (资金)", "policies (政策)", "audit_logs (审计)", "schools (学校)",
]
for i, table in enumerate(tables):
    col = i % 4
    row = i // 4
    x = 0.7 + col * 3.1
    y = 4.95 + row * 0.55

    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(2.9), Inches(0.42))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE8, 0xF0, 0xE9)
    shape.line.fill.background()
    add_text_box(slide, x, y + 0.03, 2.9, 0.35,
                 table, font_size=14, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

# Recent fixes
add_text_box(slide, 0.8, 6.2, 12, 0.4, "v1.1.0 最新修复",
             font_size=20, color=DARK_GREEN, bold=True)
recent_fixes = "启动脚本编码修复 · UserInfo菜单属性修复 · 33个冗余文档清理 · 44个脚本编码统一 · 健康检查优化"
add_text_box(slide, 0.8, 6.65, 12, 0.4, recent_fixes,
             font_size=14, color=GRAY_TEXT)

add_accent_bar(slide, 0, 7.44, 13.333, 0.06)

# ============================================================
# SLIDE 6: DEPLOYMENT
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_accent_bar(slide, 0, 0, 13.333, 0.06)

add_text_box(slide, 0.8, 0.4, 5, 0.7, "部署与交付", font_size=36, color=DARK_GREEN, bold=True)
add_accent_bar(slide, 0.8, 1.05, 2.5, 0.04)

# Deployment options
deploy_items = [
    ("Windows 安装包", "NSIS 安装程序 (.exe)\n包含 VC++ 和 WebView2 运行时\n支持 64位 / 32位 / 联合安装包\n便携版免安装选项", MID_GREEN),
    ("Linux DEB 包", "Debian/麒麟 V10 (.deb)\nARM64 架构原生支持\nsystemd 服务集成\nDocker 交叉编译构建", LIGHT_GREEN),
    ("开发环境", "一键启动脚本 (Windows .bat)\nvenv 虚拟环境隔离\n热重载开发服务器\n完整测试套件", RGBColor(0x5A, 0x7A, 0x5C)),
]

for i, (title, desc, color) in enumerate(deploy_items):
    x = 0.7 + i * 4.1
    y = 1.4
    # Card
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(3.9), Inches(3.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

    add_text_box(slide, x + 0.2, y + 0.2, 3.5, 0.5,
                 title, font_size=22, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.3, y + 0.9, 3.3, 2.3,
                 desc, font_size=15, color=RGBColor(0xEE, 0xF5, 0xEE))

# Quick start command
add_text_box(slide, 0.8, 5.3, 6, 0.5, "快速启动", font_size=24, color=DARK_GREEN, bold=True)

cmd_bg = slide.shapes.add_shape(1, Inches(0.8), Inches(5.85), Inches(7.5), Inches(0.55))
cmd_bg.fill.solid()
cmd_bg.fill.fore_color.rgb = DARK_GREEN
cmd_bg.line.fill.background()
add_text_box(slide, 1.0, 5.88, 7.1, 0.5,
             "$  scripts\\start-all.bat          # Windows 一键启动",
             font_size=16, color=WHITE, font_name='Consolas')

# System requirements
add_text_box(slide, 9.0, 5.3, 4, 0.5, "系统要求", font_size=20, color=DARK_GREEN, bold=True)
reqs = "Windows 10/11 64位 · 4GB+ RAM\nPython 3.11+ · Node.js 18+\n2GB 可用磁盘空间"
add_text_box(slide, 9.0, 5.85, 4, 1.0, reqs, font_size=13, color=GRAY_TEXT)

add_accent_bar(slide, 0, 7.44, 13.333, 0.06)

# ============================================================
# SLIDE 7: THANK YOU
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_GREEN)
add_accent_bar(slide, 0, 0, 13.333, 0.08)

add_text_box(slide, 1.5, 2.2, 10.3, 1.0,
             "感谢使用", font_size=52, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, 1.5, 3.3, 10.3, 0.8,
             "军队乡村振兴管理系统", font_size=32, color=ACCENT_GOLD, alignment=PP_ALIGN.CENTER)

add_accent_bar(slide, 5.0, 4.3, 3.3, 0.04)

contact_lines = [
    "版本 v1.1.0 · 2026年4月",
    "系统健康评分: 100/100  |  后端测试: 2812 通过  |  前端测试: 1449 通过",
    "",
    "完全离线 · 数据安全 · 高效便捷 · 生产就绪",
]
add_multi_text(slide, 2.5, 4.7, 8.3, 2.0, contact_lines,
               font_size=16, color=RGBColor(0xAA, 0xBB, 0xAA))
for p in slide.shapes[-1].text_frame.paragraphs:
    p.alignment = PP_ALIGN.CENTER

add_accent_bar(slide, 0, 7.42, 13.333, 0.08)

# Save
output_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "docs")
os.makedirs(output_dir, exist_ok=True)
output_path = os.path.join(output_dir, "军队乡村振兴管理系统_v1.1.0_演示文稿.pptx")
prs.save(output_path)
print(f"PPT saved to: {output_path}")

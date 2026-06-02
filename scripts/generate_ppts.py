#!/usr/bin/env python3
"""系统介绍PPT生成器 - 为军队乡村振兴管理系统各板块生成精美的介绍PPT"""

from concurrent.futures import ThreadPoolExecutor
from dataclasses import dataclass
from pathlib import Path
from typing import List, Union

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Theme colors
THEME_DARK_GREEN = RGBColor(0x1B, 0x43, 0x32)
THEME_GOLD = RGBColor(0xD4, 0xAF, 0x37)
THEME_LIGHT_GREEN = RGBColor(0x2D, 0x6A, 0x4F)
THEME_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
THEME_BLACK = RGBColor(0x1A, 0x1A, 0x1A)
THEME_GRAY = RGBColor(0x66, 0x66, 0x66)
THEME_LIGHT_BG = RGBColor(0xF5, 0xF7, 0xF5)

# Layout constants
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
BLANK_LAYOUT = 6  # blank slide layout index

OUTPUT_DIR = Path(__file__).resolve().parent.parent / "docs" / "PPT"


def _set_shape_color(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def _send_to_back(slide, shape):
    # Move shape to the back of the z-order.
    # Index 2 is used because positions 0/1 are reserved for group shapes.
    sp_tree = slide.shapes._spTree
    element = shape._element
    sp_tree.remove(element)
    sp_tree.insert(2, element)


def _add_title_shape(
    slide,
    text,
    left,
    top,
    width,
    height,
    font_size=44,
    bold=True,
    color=THEME_WHITE,
    align=PP_ALIGN.LEFT,
):
    tx_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tx_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Microsoft YaHei"
    p.alignment = align
    return tx_box


def _add_body_text(
    slide,
    text,
    left,
    top,
    width,
    height,
    font_size=20,
    color=THEME_BLACK,
    align=PP_ALIGN.LEFT,
    line_space=1.5,
):
    tx_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tx_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.name = "Microsoft YaHei"
    p.alignment = align
    p.line_spacing = line_space
    return tx_box


def _add_bullet_text(slide, bullets, left, top, width, height, font_size=20, color=THEME_BLACK):
    tx_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tx_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"●  {bullet}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Microsoft YaHei"
        p.space_after = Pt(12)
        p.line_spacing = 1.3
    return tx_box


def _add_header_bar(slide):
    header_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.2))
    _set_shape_color(header_bg, THEME_DARK_GREEN)
    header_bg.line.fill.background()
    _send_to_back(slide, header_bg)


def _add_gold_underline(slide):
    gold_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.05), Inches(1.5), Inches(0.05))
    _set_shape_color(gold_line, THEME_GOLD)
    gold_line.line.fill.background()


def _create_cover_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])

    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.3))
    _set_shape_color(top_bar, THEME_GOLD)
    top_bar.line.fill.background()

    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT)
    _set_shape_color(bg_shape, THEME_DARK_GREEN)
    bg_shape.line.fill.background()
    _send_to_back(slide, bg_shape)

    gold_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(3.2), Inches(2), Inches(0.05))
    _set_shape_color(gold_line, THEME_GOLD)
    gold_line.line.fill.background()

    _add_title_shape(slide, title, 1, 2.2, 10, 1, font_size=54, bold=True, color=THEME_WHITE)
    if subtitle:
        _add_title_shape(slide, subtitle, 1, 3.5, 10, 0.8, font_size=28, bold=False, color=THEME_GOLD)
    _add_body_text(slide, "军队乡村振兴管理系统 v1.1.0 · 操作培训资料", 1, 6.5, 10, 0.5, font_size=16, color=THEME_WHITE)
    _add_body_text(slide, "2026.04", 11, 6.5, 2, 0.5, font_size=16, color=THEME_GOLD)

    return slide


def _create_section_slide(prs, section_title, section_num):
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])

    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT)
    _set_shape_color(bg_shape, THEME_LIGHT_BG)
    bg_shape.line.fill.background()
    _send_to_back(slide, bg_shape)

    left_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(1.2), SLIDE_HEIGHT)
    _set_shape_color(left_bar, THEME_DARK_GREEN)
    left_bar.line.fill.background()

    _add_title_shape(slide, f"0{section_num}", 1.8, 2.5, 2, 1.5, font_size=72, bold=True, color=THEME_GOLD)
    _add_title_shape(slide, section_title, 1.8, 4, 10, 1, font_size=40, bold=True, color=THEME_DARK_GREEN)

    return slide


def _create_content_slide(prs, title, bullets, analogy=""):
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])

    _add_header_bar(slide)
    _add_title_shape(slide, title, 0.5, 0.25, 12, 0.8, font_size=32, bold=True, color=THEME_WHITE)
    _add_gold_underline(slide)
    _add_bullet_text(slide, bullets, 0.5, 1.5, 7.5, 5.5, font_size=22, color=THEME_BLACK)

    if analogy:
        analogy_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.3), Inches(1.5), Inches(4.5), Inches(5.5)
        )
        _set_shape_color(analogy_box, THEME_WHITE)
        analogy_box.line.color.rgb = THEME_GOLD
        analogy_box.line.width = Pt(3)

        analogy_title = slide.shapes.add_textbox(Inches(8.5), Inches(1.7), Inches(4.1), Inches(0.6))
        tf = analogy_title.text_frame
        p = tf.paragraphs[0]
        p.text = "💡 生活化理解"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = THEME_DARK_GREEN
        p.font.name = "Microsoft YaHei"

        analogy_body = slide.shapes.add_textbox(Inches(8.5), Inches(2.4), Inches(4.1), Inches(4.4))
        tf = analogy_body.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = analogy
        p.font.size = Pt(18)
        p.font.color.rgb = THEME_GRAY
        p.font.name = "Microsoft YaHei"
        p.line_spacing = 1.4

    return slide


def _create_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])

    _add_header_bar(slide)
    _add_title_shape(slide, title, 0.5, 0.25, 12, 0.8, font_size=32, bold=True, color=THEME_WHITE)
    _add_gold_underline(slide)

    lt = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(6), Inches(0.6))
    tf = lt.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = THEME_DARK_GREEN
    p.font.name = "Microsoft YaHei"
    _add_bullet_text(slide, left_items, 0.5, 2.0, 6, 5, font_size=20, color=THEME_BLACK)

    rt = slide.shapes.add_textbox(Inches(6.8), Inches(1.4), Inches(6), Inches(0.6))
    tf = rt.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = THEME_DARK_GREEN
    p.font.name = "Microsoft YaHei"
    _add_bullet_text(slide, right_items, 6.8, 2.0, 6, 5, font_size=20, color=THEME_BLACK)

    return slide


def _create_summary_slide(prs, title, items):
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])

    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT)
    _set_shape_color(bg_shape, THEME_DARK_GREEN)
    bg_shape.line.fill.background()
    _send_to_back(slide, bg_shape)

    _add_title_shape(slide, title, 0.8, 1.0, 12, 1, font_size=40, bold=True, color=THEME_WHITE, align=PP_ALIGN.CENTER)

    tx_box = slide.shapes.add_textbox(Inches(2), Inches(2.3), Inches(9.333), Inches(4.5))
    tf = tx_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"✓ {item}"
        p.font.size = Pt(26)
        p.font.color.rgb = THEME_WHITE
        p.font.name = "Microsoft YaHei"
        p.space_after = Pt(20)
        p.alignment = PP_ALIGN.CENTER

    return slide


@dataclass(frozen=True)
class ContentSlide:
    title: str
    bullets: List[str]
    analogy: str = ""


@dataclass(frozen=True)
class TwoColumnSlide:
    title: str
    left_title: str
    left_items: List[str]
    right_title: str
    right_items: List[str]


@dataclass(frozen=True)
class SummarySlide:
    title: str
    items: List[str]


SlideDef = Union[ContentSlide, TwoColumnSlide, SummarySlide]


@dataclass(frozen=True)
class PptSpec:
    filename: str
    title: str
    subtitle: str
    section_num: int
    slides: List[SlideDef]


PPT_SPECS = [
    PptSpec(
        filename="01-系统整体介绍.pptx",
        title="军队乡村振兴管理系统",
        subtitle="全面数字化 · 智能管理 · 一键直达",
        section_num=1,
        slides=[
            ContentSlide(
                "什么是这个系统？",
                [
                    "专为军队参与乡村振兴工作设计的信息化管理平台（v1.1.0）",
                    "覆盖帮扶村、学校、项目、资金等全流程管理，含 657 个 API 接口",
                    "支持 Windows 和国产麒麟系统（含 ARM64），87 个功能页面",
                    "完全离线运行，数据本地存储，安全可靠",
                    "内置审批工作流、RBAC 权限控制、数据级权限管理",
                ],
                "就像您手机里的'记账本+相册+日程表'的合体，但它专门用来管理军队帮扶乡村的所有工作。不需要联网，打开就能用，所有数据都存在自己的电脑里。",
            ),
            ContentSlide(
                "系统能解决什么问题？",
                [
                    "告别纸质台账，所有数据电子化、可检索",
                    "资金流向清晰透明，每一笔支出有据可查",
                    "项目进度实时跟踪，谁负责、到什么阶段一目了然",
                    "审批流程线上化，不用拿着文件到处找人签字",
                    "地图可视化展示，帮扶点位分布一图掌握",
                ],
                "以前管理乡镇工作像用一个本子记账，现在相当于有了一个智能管家：它会自动分类、自动提醒、自动汇总报表，而且永远不会把账本弄丢。",
            ),
            TwoColumnSlide(
                "系统核心技术栈",
                "后端（大脑）",
                ["Python 3.11 + FastAPI 0.109", "SQLAlchemy 2.0 + SQLite（55 模型）", "657 API 端点 + JWT 安全认证", "DDD 分层架构，完全离线运行"],
                "前端（脸面）",
                ["Vue 3.4 + TypeScript 5.9", "Element Plus 2.13（87 组件）", "Vite 7.3 + Pinia 状态管理", "ECharts 可视化 + Leaflet 地图"],
            ),
            ContentSlide(
                "系统核心功能模块",
                [
                    "帮扶村管理：建档立卡、年度数据（15 大类）、变更历史",
                    "帮扶学校管理：学校信息、助学项目、资助学生管理",
                    "帮扶项目管理：项目立项、里程碑跟踪、进度管理、验收评估",
                    "经费管理：全生命周期（申请→审批→拨付→使用→审计→决算）",
                    "审批工作流：多级审批、流程追踪、操作日志",
                    "数据管理：数据备份、质量监控、同步、加密数据包",
                    "系统管理：RBAC 权限、审计日志、系统监控、地图瓦片管理",
                ],
                "想象一个抽屉柜，每个抽屉对应一类工作：第一个抽屉装村里资料，第二个装学校资料，第三个装项目文件，第四个装账本，第五个是汇总报表。想找什么，直接拉开对应的抽屉即可。",
            ),
            SummarySlide(
                "一句话总结",
                [
                    "安全：单机离线，数据不外流",
                    "全面：村、校、项目、资金全覆盖",
                    "智能：审批、统计、预警自动化",
                    "易用：界面直观，操作简单",
                ],
            ),
        ],
    ),
    PptSpec(
        filename="02-工作台与数据看板.pptx",
        title="工作台与数据看板",
        subtitle="一屏掌控全局 · 数据一目了然",
        section_num=2,
        slides=[
            ContentSlide(
                "工作台是什么？",
                [
                    "登录后看到的第一个页面，系统的'驾驶舱'",
                    "顶部显示关键指标：帮扶村数量、项目总数、资金总额",
                    "中部展示待办事项：待审批、待完成任务",
                    "底部呈现最新动态：系统通知、工作提醒",
                ],
                "就像汽车的仪表盘，不用打开引擎盖，一眼就能看到车速、油量、水温。工作台也是这样，让您不用翻找各个菜单，就能知道当前工作的整体状况。",
            ),
            ContentSlide(
                "数据看板三大亮点",
                [
                    "多维统计：按年份、地区、项目类型自动汇总",
                    "实时图表：柱状图、饼图、折线图动态展示趋势",
                    "地图可视化：帮扶点位在地图上直观分布",
                    "一键导出：报表可导出为 Excel、Word、PDF",
                ],
                "以前领导问'今年帮扶了多少村、花了多少钱'，您可能要翻半天表格。现在点开看板，就像看电视新闻里的数据图，数字和图表一起呈现，3秒钟就能讲清楚。",
            ),
            TwoColumnSlide(
                "看板主要图表类型",
                "左侧：分类统计",
                ["饼图：资金类型占比", "柱状图：各村资金投入对比", "环形图：项目完成率分布"],
                "右侧：趋势分析",
                ["折线图：年度资金投入趋势", "面积图：项目进度累计走势", "组合图：多指标关联展示"],
            ),
            ContentSlide(
                "如何使用看板？",
                [
                    "Step 1：点击左侧菜单'数据看板'进入统计页面",
                    "Step 2：通过顶部筛选器选择年份、地区、项目类型",
                    "Step 3：鼠标悬停在图表上，查看精确数值",
                    "Step 4：点击右上角'导出'按钮，生成汇报材料",
                ],
                "操作看板就和使用手机天气APP一样简单：下拉刷新数据、滑动切换图表、点击看详情。不需要任何专业知识，会用鼠标就能上手。",
            ),
            SummarySlide(
                "核心收益",
                ["决策更快：数据实时呈现，不用等报表", "汇报更简：一键导出PPT/Excel素材", "管理更细：从村到校到项目层层穿透"],
            ),
        ],
    ),
    PptSpec(
        filename="03-帮扶村管理.pptx",
        title="帮扶村管理",
        subtitle="建档立卡 · 精准帮扶 · 动态跟踪",
        section_num=3,
        slides=[
            ContentSlide(
                "帮扶村管理管什么？",
                [
                    "基础档案：村名、地址、人口、耕地面积、联系方式",
                    "年度数据：人均收入、产业收入、基础设施变化",
                    "帮扶记录：历年投入的项目、资金、人员",
                    "变更历史：村干部变动、脱贫摘帽等重要节点",
                ],
                "就像给每个帮扶村建立了一份'病历档案'，从第一次挂牌帮扶开始，所有的检查结果（年度数据）、开的药方（项目）、花的医药费（资金）都清清楚楚记录在案。",
            ),
            ContentSlide(
                "关键操作流程",
                [
                    "新增村庄：填写基本信息 → 上传照片 → 保存建档",
                    "录入年度数据：选择年份 → 填写经济指标 → 提交审核",
                    "查看详情：点击村名 → 查看历年数据和关联项目",
                    "编辑更新：发现信息有误 → 点击编辑 → 修改保存",
                ],
                "和在微信里添加一位新联系人非常像：输入姓名电话（基础信息），备注分组（所属地区），然后每年更新一下对方的近况（年度数据），保持档案不过期。",
            ),
            ContentSlide(
                "年度数据有哪些指标？",
                [
                    "经济指标：人均纯收入、集体经济收入、产业收入",
                    "人口指标：总人口、劳动力人数、脱贫人口数",
                    "基础设施：道路硬化里程、饮水安全户数、网络覆盖率",
                    "教育医疗：适龄儿童入学率、村卫生室覆盖率",
                ],
                "这些指标就像体检单上的血压、血糖、心率。通过每年对比，可以清楚看出这个村的'健康状况'是在好转还是退步，帮扶效果是否明显。",
            ),
            SummarySlide(
                "管理要点",
                ["一户一档：每个村都有独立完整的电子档案", "逐年对比：历史数据自动形成趋势线", "关联穿透：从村庄可直接进入项目和资金"],
            ),
        ],
    ),
    PptSpec(
        filename="04-帮扶学校管理.pptx",
        title="帮扶学校管理",
        subtitle="助学兴教 · 点亮希望",
        section_num=4,
        slides=[
            ContentSlide(
                "帮扶学校管理功能",
                [
                    "学校档案：名称、类型（小学/初中/高中）、师生人数",
                    "帮扶项目：援建图书室、硬化操场、捐赠教学设备",
                    "资助学生：困难学生名单、资助金额、资助学期",
                    "成效跟踪：学生成绩变化、辍学率变化、升学率",
                ],
                "这个功能就像一本'爱心助学账本'。哪所学校得到了什么帮助、哪些学生收到了资助、效果怎么样，翻开账本就能找到明细。",
            ),
            ContentSlide(
                "如何添加一所新学校？",
                [
                    "Step 1：进入'帮扶学校'菜单，点击'新增学校'",
                    "Step 2：填写学校名称、地址、办学类型、师生人数",
                    "Step 3：上传校门照片和校园环境照片",
                    "Step 4：保存后自动关联到对应的帮扶村",
                ],
                "和在通讯录里添加一个新单位很像：写下名字和地址，备注一下有什么合作往来，之后双方的每一次互动都可以在这里追加记录。",
            ),
            ContentSlide(
                "资助学生管理怎么做？",
                [
                    "建立名单：录入受助学生姓名、年级、家庭情况",
                    "发放记录：记录每笔资助金额、发放时间、发放人",
                    "学期跟踪：每学期更新学生的学习情况和在校状态",
                    "隐私保护：敏感信息仅限授权人员查看",
                ],
                "每一笔助学金的流向都像银行转账一样有记录。家长、学校、部队三方都能看到自己关心的那部分信息，确保助学金发到真正需要的孩子手里。",
            ),
            SummarySlide(
                "核心价值",
                ["档案完整：从学校到学生全覆盖", "跟踪到位：从资助到成效可追溯", "关爱传递：让每一分帮扶都落到实处"],
            ),
        ],
    ),
    PptSpec(
        filename="05-帮扶项目管理.pptx",
        title="帮扶项目管理",
        subtitle="全程跟踪 · 闭环管理",
        section_num=5,
        slides=[
            ContentSlide(
                "项目管理管什么？",
                [
                    "项目立项：名称、预算、负责人、计划起止时间",
                    "进度跟踪：定期更新完成百分比、当前阶段、存在问题",
                    "资金管理：项目关联经费的申请、拨付、使用",
                    "验收评估：项目完成后上传验收材料、打分评价",
                ],
                "想象您要装修一套房子：先要立项（确定预算和工期），施工中每周汇报进度（刷墙完成了60%），最后验收签字（确认质量合格）。项目管理就是这样一套流程。",
            ),
            ContentSlide(
                "项目生命周期（5个阶段）",
                [
                    "1. 论证立项：提出申请，阐明项目必要性和预期效果",
                    "2. 汇总审核：上级审批预算，确认项目可行性",
                    "3. 计划下达：批准预算，资金拨付到位",
                    "4. 组织实施：进场施工，定期汇报进度",
                    "5. 核查决算：完工验收，审计决算，总结归档",
                ],
                "就像种一棵树：先选树苗（立项），再挖坑施肥（拨款），然后栽树浇水（施工），最后检查成活情况（验收）。每个环节都不能跳过，否则树可能长不好。",
            ),
            TwoColumnSlide(
                "项目进度怎么看？",
                "电脑端",
                ["列表视图：所有项目一览表", "详情页：单个项目的完整记录", "甘特图：时间计划和实际对比"],
                "移动端（如有）",
                ["拍照上传：现场施工照片", "快速填报：简版进度更新", "消息推送：待办审批提醒"],
            ),
            SummarySlide(
                "管好项目的秘诀",
                ["有计划：立项时明确目标和预算", "有跟踪：每周更新进度不拖延", "有闭环：完工必验收，花钱必对账"],
            ),
        ],
    ),
    PptSpec(
        filename="06-经费管理.pptx",
        title="经费管理",
        subtitle="精打细算 · 阳光透明",
        section_num=6,
        slides=[
            ContentSlide(
                "经费管理为什么重要？",
                [
                    "资金是帮扶工作的'血液'，必须管好每一分钱",
                    "系统实现预算、拨付、使用、决算全生命周期管理",
                    "每笔资金都与项目和村庄绑定，杜绝虚列支出",
                    "异常自动预警：超预算、拨付滞后、执行率过低都会提醒",
                ],
                "就像家里的理财APP：您设定了每月伙食费2000元（预算），每次买菜支出都记账（拨付），月底系统自动告诉您花超了还是剩了（决算），还能发现哪笔支出异常（预警）。",
            ),
            ContentSlide(
                "经费全生命周期管理（6 大阶段）",
                [
                    "申请阶段：在线提交经费申请，关联项目和村庄",
                    "审批阶段：多级审批（通过/驳回/转审），附带审批意见",
                    "拨付阶段：生成划转凭证，确认资金到账",
                    "使用阶段：合同管理、支付记录、WBS 编码关联",
                    "完成阶段：决算报告、资产清查、绩效评分",
                    "审计阶段：审计意见、绩效等级（优秀/良好/合格/不合格）",
                ],
                "想象一笔经费是一块蛋糕：先确定蛋糕做多大（预算），然后切一块给家人（拨付），每次有人吃了蛋糕都要登记（支出），最后看蛋糕够不够分（决算）。",
            ),
            ContentSlide(
                "异常预警机制",
                [
                    "预算执行率偏差：实际支出和计划差距超过设定阈值时预警",
                    "拨付进度滞后：项目已过半但资金未拨付，亮黄灯提醒",
                    "重复支付风险：同一笔款项多次申请时自动拦截",
                    "合同付款异常：合同约定分三期，若第二期超期未付则预警",
                ],
                "这就像您手机银行的短信提醒：大额支出会收到短信，余额不足会收到提示。系统的预警功能就是帮您当好'财务管家'，让风险早发现、早处理。",
            ),
            SummarySlide(
                "管好经费的4个关键词",
                ["预算先行：没有预算不花钱", "据实支出：每笔花销有凭据", "审批留痕：线上审批可追溯", "预警及时：异常早发现早处理"],
            ),
        ],
    ),
    PptSpec(
        filename="07-审批工作流.pptx",
        title="审批工作流",
        subtitle="线上流转 · 高效透明",
        section_num=7,
        slides=[
            ContentSlide(
                "什么是审批工作流？",
                [
                    "把原来'拿着文件找领导签字'的流程搬到了线上",
                    "申请人提交 → 审批人收到待办 → 点击通过/驳回",
                    "全程留痕：谁审批的、什么时间、什么意见都自动记录",
                    "审批结果实时通知申请人，不用反复打电话询问",
                ],
                "就像您在网上申请信用卡：填表提交后，银行后台审核，然后App里直接显示'已通过'或'已驳回'。审批工作流就是让单位的内部审批也能像银行那样高效透明。",
            ),
            ContentSlide(
                "常见的审批场景",
                [
                    "经费申请审批：用款人申请 → 分管领导审批 → 财务复核",
                    "项目立项审批：项目负责人提交 → 机关审批 → 领导签批",
                    "用户注册审批：新用户注册后，管理员审核激活账号",
                    "数据变更审批：重要数据修改需经审批后才生效",
                ],
                "企事业单位常见的请假条、报销单、用印申请，本质都是'先申请再审批'的流程。系统把这些流程标准化，申请人不用跑腿，审批人不会漏审。",
            ),
            ContentSlide(
                "审批中心怎么用？",
                [
                    "我的申请：查看自己提交的所有申请及当前状态",
                    "待办审批：显示所有需要我审批的事项，一目了然",
                    "审批历史：查看已审批事项的完整记录和审批意见",
                    "流程追踪：像查快递一样，实时看审批进行到哪一步",
                ],
                "和网购平台的'我的订单'非常像：待付款（待提交）、待发货（待审批）、已发货（已通过）、已退款（已驳回）。每个状态都很清晰，点击还能看物流（流程详情）。",
            ),
            SummarySlide(
                "使用审批流的好处",
                ["省时：不用拿着文件到处找人", "留痕：全程记录，责任清晰", "透明：申请人随时可查进度", "规范：审批标准统一，不走人情"],
            ),
        ],
    ),
    PptSpec(
        filename="08-权限与安全管理.pptx",
        title="权限与安全管理",
        subtitle="分级管控 · 安全无忧",
        section_num=8,
        slides=[
            ContentSlide(
                "权限管理的核心逻辑",
                [
                    "角色决定权限：管理员、普通用户、只读用户权限不同",
                    "数据范围控制：有人看全部，有人只看自己村的数据",
                    "菜单级控制：不同角色看到的系统菜单不一样",
                    "按钮级控制：同一页面，有人能编辑，有人只能查看",
                ],
                "就像公司大楼的门禁卡：总经理的卡能开所有门（管理员），部门经理能开自己部门的门（数据范围），实习生只能进公共区（只读权限）。每个人都只能在自己权限范围内活动。",
            ),
            ContentSlide(
                "系统有哪些角色？",
                [
                    "超级管理员（super_admin）：拥有最高权限，可管理用户、配置系统",
                    "管理员（admin）：用户管理、模块管理、数据导出",
                    "管理者（manager）：本组织及下级业务模块 CRUD、审批、报表",
                    "审批领导（approval_leader）：本组织审批操作、查看、导出",
                    "查看者（viewer）：仅自己创建的数据，只读查看",
                ],
                "这和医院的医护分工类似：院长管全院（管理员），科主任管科室（业务管理员），护士负责具体病人记录（操作员），病人家属只能看公示栏（查看者）。",
            ),
            ContentSlide(
                "安全保障措施",
                [
                    "单机离线：系统不联网，从根本上杜绝网络攻击",
                    "密码加密：用户密码采用 bcrypt 加密（10 轮），即使数据库泄露也无法破解",
                    "JWT 令牌：登录后发放安全令牌（8h 有效），支持 Token 黑名单和双因素认证",
                    "四级数据权限：all（全部）/ org_children（本组织及下级）/ org（本组织）/ self（仅自己）",
                    "操作审计：全操作审计日志 + 安全事件告警 + API 访问日志",
                    "机器码绑定：可绑定特定电脑，防止软件随意复制使用",
                ],
                "保险箱为什么安全？因为它有厚厚的钢板（离线）、复杂的密码锁（加密）、开锁记录（审计）。系统的安全措施就像多套保险箱的组合，层层保护您的数据。",
            ),
            SummarySlide(
                "安全理念",
                ["最小权限：只给工作必需的权限", "全程留痕：关键操作都有记录", "本地优先：数据不出本地电脑"],
            ),
        ],
    ),
    PptSpec(
        filename="09-离线地图与数据管理.pptx",
        title="离线地图与数据管理",
        subtitle="一图纵览 · 数据随行",
        section_num=9,
        slides=[
            ContentSlide(
                "离线地图有什么作用？",
                [
                    "不需要网络就能查看地图和帮扶点位分布",
                    "在地图上标注帮扶村、学校、项目的位置",
                    "支持地图搜索、距离测算、区域覆盖分析",
                    "可预先下载瓦片到本地，打开就能流畅浏览",
                ],
                "就像手机里下载好的离线导航地图。即使到了没有信号的山区，也能打开地图看自己在哪、附近有哪些帮扶点。不用担心网络问题影响工作。",
            ),
            ContentSlide(
                "数据导入导出功能",
                [
                    "批量导入：村、校、项目等数据可通过 Excel 批量导入",
                    "数据导出：一键导出报表为 Excel / Word / PDF",
                    "加密数据包：敏感数据可加密打包，安全转移",
                    "数据备份：支持手动备份和自动定时备份",
                ],
                "就像手机的'通讯录备份'功能：可以把联系人导出到电脑上保存（数据导出），也可以把电脑上的联系人列表批量导入手机（数据导入），还能加密压缩后传给同事。",
            ),
            ContentSlide(
                "数据质量管理",
                [
                    "完整性检查：必填项缺失时自动标红提醒",
                    "一致性校验：同一数据在不同模块中是否一致",
                    "重复检测：导入时自动识别重复记录",
                    "历史版本：重要数据变更保留历史版本，可随时回退",
                ],
                "就像考试阅卷时的检查：看名字写了没有（完整性）、选择题涂卡是否和答案一致（一致性）、有没有两份同名试卷（重复检测）。系统会自动帮您做这些检查。",
            ),
            SummarySlide(
                "核心价值",
                ["无网可用：离线地图保障山区工作", "数据互通：导入导出让协作更便捷", "质量可控：自动检查减少人为错误"],
            ),
        ],
    ),
    PptSpec(
        filename="10-系统部署与维护.pptx",
        title="系统部署与维护",
        subtitle="三步上手 · 轻松运维",
        section_num=10,
        slides=[
            ContentSlide(
                "支持哪些平台？",
                [
                    "Windows 10/11（x64）：最常见的办公电脑环境",
                    "麒麟 V10（x64）：国产操作系统标准版本",
                    "麒麟 V10（ARM64）：飞腾/鲲鹏芯片的国产电脑",
                    "全部是单机离线部署，不依赖外网",
                ],
                "就像同一款软件同时出了 Windows 版、Mac 版和国产手机版。无论您用哪种电脑，都有一键安装包可以直接运行。",
            ),
            ContentSlide(
                "三步完成安装",
                [
                    "Step 1：下载安装包（Windows 是 .exe，麒麟是 .deb）",
                    "Step 2：双击安装，选择安装目录，点击'下一步'",
                    "Step 3：安装完成后点击桌面图标，即可启动系统",
                ],
                "安装过程和安装微信、WPS 等日常软件完全一样：双击安装程序、点几次下一步、等待进度条跑完。没有任何技术门槛，普通文员也能独立完成。",
            ),
            TwoColumnSlide(
                "日常维护做什么？",
                "数据安全类",
                ["定期备份数据库", "备份文件异地保存", "清理过期日志文件"],
                "系统优化类",
                ["定期检查和升级版本", "监控系统运行状态", "及时清理临时文件"],
            ),
            ContentSlide(
                "常见问题和解决办法",
                [
                    "启动卡慢？检查磁盘空间是否充足，清理旧日志",
                    "忘记密码？联系系统管理员重置密码",
                    "数据丢失？通过备份文件恢复数据库",
                    "界面显示异常？尝试清除浏览器缓存后重试",
                ],
                "就像汽车需要定期保养（换机油、检查轮胎），系统也需要定期做简单维护。但这些工作都很简单，不需要专业技术人员，按手册操作即可。",
            ),
            SummarySlide(
                "记住这3点",
                ["有安装包：Windows / 麒麟 一键安装", "有备份机制：数据可手动+自动双重备份", "有恢复方案：出问题可按手册快速恢复"],
            ),
        ],
    ),
]


def _build_slide(prs, slide_def: SlideDef):
    if isinstance(slide_def, ContentSlide):
        _create_content_slide(prs, slide_def.title, slide_def.bullets, slide_def.analogy)
    elif isinstance(slide_def, TwoColumnSlide):
        _create_two_column_slide(
            prs,
            slide_def.title,
            slide_def.left_title,
            slide_def.left_items,
            slide_def.right_title,
            slide_def.right_items,
        )
    elif isinstance(slide_def, SummarySlide):
        _create_summary_slide(prs, slide_def.title, slide_def.items)


def generate_ppt(spec: PptSpec) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    _create_cover_slide(prs, spec.title, spec.subtitle)
    _create_section_slide(prs, spec.title, spec.section_num)

    for slide_def in spec.slides:
        _build_slide(prs, slide_def)

    out_path = OUTPUT_DIR / spec.filename
    prs.save(str(out_path))
    print(f"Generated {spec.filename}")


def main():
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    with ThreadPoolExecutor(max_workers=4) as executor:
        futures = [executor.submit(generate_ppt, spec) for spec in PPT_SPECS]
        for future in futures:
            future.result()

    print("All PPTs generated successfully!")


if __name__ == "__main__":
    main()

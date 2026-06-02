"""
生成军民融合乡村振兴系统 V1.0.4 更新说明 PPT
"""
import sys
import os

# 添加项目根目录到路径
sys.path.insert(0, os.path.dirname(__file__))

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print("正在安装 python-pptx...")
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# 定义颜色方案
COLOR_PRIMARY = RGBColor(41, 128, 185)  # 蓝色
COLOR_SUCCESS = RGBColor(39, 174, 96)   # 绿色
COLOR_WARNING = RGBColor(243, 156, 18)  # 橙色
COLOR_DANGER = RGBColor(231, 76, 60)    # 红色
COLOR_DARK = RGBColor(44, 62, 80)       # 深灰色
COLOR_LIGHT = RGBColor(236, 240, 241)   # 浅灰色

def add_title_slide(prs, title, subtitle):
    """添加标题幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

    # 添加背景色
    background = slide.shapes.add_shape(
        1,  # 矩形
        0, 0,
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLOR_PRIMARY
    background.line.fill.background()

    # 添加标题
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5),
        Inches(8), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER

    # 添加副标题
    subtitle_box = slide.shapes.add_textbox(
        Inches(1), Inches(4.5),
        Inches(8), Inches(1)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
    subtitle_para.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_items):
    """添加内容幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 添加标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5),
        Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_PRIMARY

    # 添加内容
    content_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.5),
        Inches(8.4), Inches(5.5)
    )
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    for item in content_items:
        p = text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.space_before = Pt(12)
        p.level = 0

def add_table_slide(prs, title, headers, rows):
    """添加表格幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 添加标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5),
        Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_PRIMARY

    # 添加表格
    rows_count = len(rows) + 1
    cols_count = len(headers)

    table = slide.shapes.add_table(
        rows_count, cols_count,
        Inches(0.8), Inches(1.8),
        Inches(8.4), Inches(4.5)
    ).table

    # 设置表头
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_PRIMARY
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(14)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(255, 255, 255)

    # 填充数据
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_data)
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(12)

# 幻灯片 1: 封面
add_title_slide(
    prs,
    "军民融合乡村振兴管理系统",
    "V1.0.4 重大安全与性能优化更新"
)

# 幻灯片 2: 更新概览
add_content_slide(
    prs,
    "更新概览",
    [
        "✅ 密钥安全管理 - 企业级加密存储",
        "✅ Token 黑名单持久化 - 双层缓存架构",
        "✅ 权限过滤 SQL 优化 - 性能提升 70-90%",
        "✅ 缓存策略优化 - 命中率提升至 85%+",
        "✅ 慢查询优化 - 阈值降低至 200ms",
        "✅ CSRF 保护启用 - 增强安全防护",
        "",
        "📅 发布日期: 2026-03-10",
        "📦 版本号: V1.0.4",
        "🎯 更新类型: 重大安全与性能优化"
    ]
)

# 幻灯片 3: 安全性增强 - 密钥管理
add_content_slide(
    prs,
    "安全性增强 - 密钥加密存储",
    [
        "🔒 问题: 敏感密钥以明文存储，存在泄露风险",
        "",
        "✅ 解决方案:",
        "  • 使用 Fernet (AES-256) 对称加密",
        "  • 采用 PBKDF2 密钥派生（100,000 次迭代）",
        "  • 主密钥文件权限 0600",
        "  • 支持环境变量回退机制",
        "",
        "📈 安全提升:",
        "  • 消除明文密钥泄露风险",
        "  • 符合企业级安全标准",
        "  • 支持密钥轮换机制"
    ]
)

# 幻灯片 4: 安全性增强 - Token 黑名单
add_content_slide(
    prs,
    "安全性增强 - Token 黑名单持久化",
    [
        "🔒 问题: Token 黑名单仅在内存，重启后丢失",
        "",
        "✅ 解决方案:",
        "  • SQLite + 缓存双层架构",
        "  • Token 黑名单持久化到数据库",
        "  • 内存缓存加速查询（命中率 >95%）",
        "  • 自动清理过期 Token",
        "",
        "📈 功能特性:",
        "  • 支持服务重启后保留",
        "  • 支持多种撤销原因",
        "  • 双层缓存架构"
    ]
)

# 幻灯片 5: 性能优化 - 权限过滤
add_content_slide(
    prs,
    "性能优化 - 权限过滤 SQL 优化",
    [
        "⚡ 问题: 应用层权限过滤，N+1 查询问题",
        "",
        "✅ 解决方案:",
        "  • 使用递归 CTE 查询组织层级",
        "  • 在 SQL 层面实现权限过滤",
        "  • 单次 SQL 查询完成过滤",
        "  • 组织层级缓存（5 分钟 TTL）",
        "",
        "📈 性能提升:",
        "  • 查询时间减少 70-90%",
        "  • 消除 N+1 查询问题",
        "  • 单次 SQL 完成过滤"
    ]
)

# 幻灯片 6: 性能优化 - 缓存策略
add_content_slide(
    prs,
    "性能优化 - 缓存策略优化",
    [
        "⚡ 问题: 缓存命中率不足（约 60%）",
        "",
        "✅ 解决方案:",
        "  • 缓存预热服务（启动时预热热点数据）",
        "  • 缓存监控工具（记录命中率）",
        "  • 优化 TTL 策略（动态调整）",
        "  • 智能缓存失效（级联失效）",
        "",
        "📈 性能提升:",
        "  • 缓存命中率: 60% → 85%+",
        "  • 数据库查询减少 40%+",
        "  • 响应时间降低 30%+"
    ]
)

# 幻灯片 7: 性能指标对比
add_table_slide(
    prs,
    "性能指标对比",
    ["性能指标", "优化前", "优化后", "提升幅度"],
    [
        ["权限过滤查询", "基准", "减少 70-90%", "⬇️ 70-90%"],
        ["缓存命中率", "~60%", "~85%+", "⬆️ 25%"],
        ["数据库查询数", "基准", "减少 40%+", "⬇️ 40%"],
        ["平均响应时间", "基准", "降低 30%+", "⬇️ 30%"],
        ["慢查询阈值", "500ms", "200ms", "⬇️ 60%"]
    ]
)

# 幻灯片 8: 安全指标对比
add_table_slide(
    prs,
    "安全指标对比",
    ["安全问题", "优化前", "优化后", "状态"],
    [
        ["密钥存储", "明文存储", "AES-256 加密", "✅ 已解决"],
        ["Token 黑名单", "内存缓存", "持久化存储", "✅ 已解决"],
        ["CSRF 保护", "未启用", "已启用", "✅ 已解决"],
        ["日志脱敏", "部分脱敏", "完全脱敏", "✅ 已增强"]
    ]
)

# 幻灯片 9: 架构改进
add_content_slide(
    prs,
    "架构改进",
    [
        "🏗️ 双层缓存架构:",
        "  • 内存缓存（热点数据，快速访问）",
        "  • 磁盘缓存（持久化存储，100MB 限制）",
        "",
        "🏗️ SQL 层权限过滤:",
        "  • 递归 CTE 查询组织层级",
        "  • 组织层级缓存",
        "  • 消除 N+1 查询",
        "",
        "🏗️ Token 黑名单持久化:",
        "  • SQLite 持久化存储",
        "  • 双层缓存加速",
        "  • 自动清理机制"
    ]
)

# 幻灯片 10: 新增文件
add_content_slide(
    prs,
    "新增文件清单",
    [
        "📁 密钥管理 (3 个文件):",
        "  • secret_migration.py - 密钥迁移工具",
        "  • encrypted_config.json - 加密配置",
        "  • master.key - 主密钥",
        "",
        "📁 Token 黑名单 (2 个文件):",
        "  • token_blacklist.py - 数据模型",
        "  • token_blacklist_service.py - 服务层",
        "",
        "📁 权限过滤 (2 个文件):",
        "  • permission_filter.py - SQL 层过滤",
        "  • user_organization.py - 用户组织模型"
    ]
)

# 幻灯片 11: 数据库变更
add_content_slide(
    prs,
    "数据库变更",
    [
        "🗄️ 新增数据表:",
        "",
        "1. token_blacklist - Token 黑名单表",
        "   • 5 个字段（id, token_jti, user_id, ...）",
        "   • 2 个索引（user_time, expires）",
        "",
        "2. user_organizations - 用户-组织关联表",
        "   • 6 个字段（id, user_id, organization_id, ...）",
        "   • 3 个索引（user_org, user, org）",
        "   • 1 个唯一约束",
        "",
        "📝 迁移脚本:",
        "   • migrate_token_blacklist.py",
        "   • migrate_user_organizations.py"
    ]
)

# 幻灯片 12: 升级指南
add_content_slide(
    prs,
    "升级指南",
    [
        "📋 升级步骤:",
        "",
        "1️⃣ 备份数据（数据库 + 配置文件）",
        "2️⃣ 更新代码（git pull origin main）",
        "3️⃣ 运行数据库迁移",
        "   python migrate_token_blacklist.py",
        "   python migrate_user_organizations.py",
        "4️⃣ 迁移密钥（可选）",
        "   python -m backend.app.utils.secret_migration",
        "5️⃣ 更新配置（慢查询阈值、CSRF）",
        "6️⃣ 重启应用",
        "7️⃣ 验证升级（运行测试脚本）"
    ]
)

# 幻灯片 13: 兼容性说明
add_content_slide(
    prs,
    "兼容性说明",
    [
        "✅ 完全向后兼容",
        "",
        "• 所有新功能都是可选的",
        "• 不启用加密密钥时，使用环境变量",
        "• Token 黑名单自动迁移",
        "• 权限过滤逻辑保持兼容",
        "",
        "💻 系统要求:",
        "  • Python 3.11+",
        "  • SQLite 3.35+",
        "  • Node.js 18+",
        "  • 磁盘空间: 2GB+",
        "  • 内存: 4GB+"
    ]
)

# 幻灯片 14: 技术亮点
add_content_slide(
    prs,
    "技术亮点",
    [
        "✨ 使用递归 CTE 优化组织层级查询",
        "",
        "✨ 实现双层缓存架构（内存 + 磁盘）",
        "",
        "✨ 使用 Fernet + PBKDF2 加密密钥",
        "",
        "✨ 实现智能缓存预热和失效",
        "",
        "✨ 完整的缓存监控体系",
        "",
        "✨ 参数脱敏防止敏感信息泄露"
    ]
)

# 幻灯片 15: 总结
add_content_slide(
    prs,
    "总结",
    [
        "🎯 本次更新成果:",
        "",
        "✅ 6 个核心优化模块全部完成",
        "✅ 安全性大幅提升（消除 3 个高风险问题）",
        "✅ 性能显著优化（响应时间降低 30%+）",
        "✅ 架构全面改进（双层缓存、SQL 层过滤）",
        "✅ 完整的测试验证",
        "✅ 详细的文档说明",
        "",
        "🚀 系统安全性和性能得到全面提升",
        "   为后续业务发展奠定坚实技术基础"
    ]
)

# 幻灯片 16: 致谢
add_title_slide(
    prs,
    "感谢使用",
    "军民融合乡村振兴管理系统 V1.0.4"
)

# 保存演示文稿
output_file = "军民融合乡村振兴系统_V1.0.4_更新说明.pptx"
prs.save(output_file)

print(f"[OK] PPT generated successfully: {output_file}")
print(f"[INFO] Total slides: 16")
print(f"[INFO] File location: {os.path.abspath(output_file)}")
